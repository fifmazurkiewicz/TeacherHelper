"""Konwersja treści plików do formatów eksportowych (PDF, DOCX, TXT, PPTX)."""
from __future__ import annotations

import io
import platform
import re
from pathlib import Path

SUPPORTED_FORMATS = ("txt", "pdf", "docx", "pptx")


def text_to_txt(text: str) -> bytes:
    return text.encode("utf-8")


def _wrap_line_for_pdf(line: str, max_chars: int = 92) -> list[str]:
    """Dzieli linię na fragmenty mieszczące się w PDF (fpdf wywala się na jednym „słowie” szerzej niż strona)."""
    line = line.replace("\r", "").replace("\t", "    ")
    if len(line) <= max_chars:
        return [line] if line else [" "]
    out: list[str] = []
    i = 0
    n = len(line)
    while i < n:
        end = min(i + max_chars, n)
        if end < n:
            sp = line.rfind(" ", i + 1, end)
            if sp > i + 8:
                end = sp + 1
        chunk = line[i:end].rstrip()
        out.append(chunk if chunk else " ")
        if end <= i:
            end = i + 1
        i = end
    return out or [" "]


def text_to_pdf(text: str, title: str = "") -> bytes:
    from fpdf import FPDF
    from fpdf.enums import XPos, YPos

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_margins(14, 14, 14)
    pdf.add_page()

    font_path = _find_unicode_font()
    if font_path:
        pdf.add_font("Unicode", "", font_path)
        pdf.set_font("Unicode", size=11)
    else:
        pdf.set_font("Helvetica", size=11)

    if title:
        pdf.set_font_size(16)
        pdf.set_x(pdf.l_margin)
        pdf.multi_cell(0, 10, title, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font_size(11)
        pdf.ln(3)

    for raw in text.split("\n"):
        for seg in _wrap_line_for_pdf(raw):
            pdf.set_x(pdf.l_margin)
            pdf.multi_cell(0, 6, seg if seg.strip() else " ", new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    return bytes(pdf.output())


def text_to_docx(text: str, title: str = "") -> bytes:
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    if title:
        doc.add_heading(title, level=1)
    for para_text in text.split("\n"):
        p = doc.add_paragraph(para_text)
        for run in p.runs:
            run.font.size = Pt(11)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def text_to_pptx(text: str, title: str = "") -> bytes:
    """Generuje prezentację PPTX z tekstu markdown-like (nagłówki = slajdy)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = _parse_slides(text, title)

    for slide_title, bullets in slides_data:
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = slide_title

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        for i, bullet in enumerate(bullets):
            if i == 0:
                tf.paragraphs[0].text = bullet
                tf.paragraphs[0].font.size = Pt(18)
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)

    if not slides_data:
        layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title or "Prezentacja"
        if slide.placeholders[1]:
            slide.placeholders[1].text = text[:500]

    from teacher_helper.infrastructure.presentation_spec import apply_colorful_theme_to_presentation

    apply_colorful_theme_to_presentation(prs)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _parse_slides(text: str, fallback_title: str) -> list[tuple[str, list[str]]]:
    """Parsuje tekst na slajdy: nagłówki markdown (#/##) → tytuły, reszta → punktory."""
    slides: list[tuple[str, list[str]]] = []
    current_title = ""
    current_bullets: list[str] = []

    for line in text.split("\n"):
        stripped = line.strip()
        heading_match = re.match(r"^#{1,3}\s+(.+)$", stripped)
        if heading_match:
            if current_title or current_bullets:
                slides.append((current_title or fallback_title, current_bullets))
            current_title = heading_match.group(1).strip()
            current_bullets = []
        elif stripped.startswith(("- ", "* ", "• ")):
            current_bullets.append(stripped.lstrip("-*• ").strip())
        elif re.match(r"^\d+\.\s+", stripped):
            current_bullets.append(re.sub(r"^\d+\.\s+", "", stripped).strip())
        elif stripped:
            current_bullets.append(stripped)

    if current_title or current_bullets:
        slides.append((current_title or fallback_title, current_bullets))

    return slides


def convert_text(text: str, target_format: str, title: str = "") -> tuple[bytes, str]:
    """Konwertuje tekst do docelowego formatu. Zwraca (bytes, mime_type)."""
    fmt = target_format.lower().strip().lstrip(".")
    if fmt == "txt":
        return text_to_txt(text), "text/plain; charset=utf-8"
    if fmt == "pdf":
        return text_to_pdf(text, title), "application/pdf"
    if fmt == "docx":
        return text_to_docx(text, title), (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
    if fmt == "pptx":
        return text_to_pptx(text, title), (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    raise ValueError(f"Nieobsługiwany format eksportu: {fmt}. Dostępne: {', '.join(SUPPORTED_FORMATS)}")


def _find_unicode_font() -> str | None:
    """Szuka systemowej czcionki TTF ze wsparciem dla polskich znaków."""
    system = platform.system()
    candidates: list[Path] = []
    if system == "Windows":
        base = Path("C:/Windows/Fonts")
        candidates = [base / "arial.ttf", base / "calibri.ttf", base / "segoeui.ttf"]
    elif system == "Darwin":
        candidates = [
            Path("/System/Library/Fonts/Supplemental/Arial.ttf"),
            Path("/Library/Fonts/Arial.ttf"),
        ]
    else:
        candidates = [
            Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
            Path("/usr/share/fonts/TTF/DejaVuSans.ttf"),
            Path("/usr/share/fonts/dejavu-sans-fonts/DejaVuSans.ttf"),
        ]
    for p in candidates:
        if p.exists():
            return str(p)
    return None
