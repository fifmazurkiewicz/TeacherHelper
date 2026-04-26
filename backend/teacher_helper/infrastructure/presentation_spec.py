"""Struktura prezentacji (JSON) → PPTX, odczyt PPTX/JSON, tekst do indeksu."""
from __future__ import annotations

import io
import json
import re
from typing import Any

# --- Normalizacja specyfikacji z LLM / z pliku JSON ---

# Domyślna paleta (gdy brak / częściowy `theme` w specie)
_DEFAULT_THEME_RGB: dict[str, tuple[int, int, int]] = {
    "background": (0x1A, 0x24, 0x3A),
    "title": (0x7E, 0xC8, 0xF0),
    "body": (0xE8, 0xF0, 0xF8),
    "muted": (0x9A, 0xB4, 0xCC),
}
# Gdy tło slajdu jest jasne, a model zwróci jasne litery — podmiana na czytelne ciemne
_DEFAULT_DARK_ON_LIGHT: dict[str, tuple[int, int, int]] = {
    "title": (0x0D, 0x1B, 0x2A),
    "body": (0x1B, 0x26, 0x3B),
    "muted": (0x41, 0x5A, 0x77),
}
_THEME_JSON_KEYS: tuple[str, ...] = ("background", "title", "body", "muted")
_BG_LUMINANCE_LIGHT: float = 0.55  # powyżej: jasne tło → ciemny tekst z zapasu


def _parse_hex_rgb(s: str) -> tuple[int, int, int] | None:
    t = s.strip().lstrip("#")
    if len(t) == 3 and all(c in "0123456789abcdefABCDEF" for c in t):
        t = "".join(c * 2 for c in t)
    if len(t) != 6 or not all(c in "0123456789abcdefABCDEF" for c in t):
        return None
    try:
        return (int(t[0:2], 16), int(t[2:4], 16), int(t[4:6], 16))
    except ValueError:
        return None


def _luminance_srgb(r: int, g: int, b: int) -> float:
    def ch(x: int) -> float:
        c = x / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

    return 0.2126 * ch(r) + 0.7152 * ch(g) + 0.0722 * ch(b)


def _contrast_ratio_rgb(fg: tuple[int, int, int], bg: tuple[int, int, int]) -> float:
    """Współczynnik kontrastu WCAG 2.1 (większa / mniejsza jasność + 0,05)."""
    lf, lb = _luminance_srgb(*fg), _luminance_srgb(*bg)
    a, b = (lf, lb) if lf >= lb else (lb, lf)
    if b < 0:
        return 21.0
    return (a + 0.05) / (b + 0.05)


def _mend_theme_contrast(raw: dict[str, str] | None) -> dict[str, str] | None:
    """
    Gdy model zwróci tło i tekst zbyt zbliżone (np. ciemny body na ciemnym tle),
    podmienia kolory tytułu/treści/akcentu na czytelną parę: jasne litery na ciemnym tle
    albo ciemne na jasnym (wg jasności tła).
    """
    if not raw:
        return None
    resolved: dict[str, tuple[int, int, int]] = {}
    for k in _THEME_JSON_KEYS:
        t = _parse_hex_rgb(str(raw.get(k) or ""))
        resolved[k] = t if t else _DEFAULT_THEME_RGB[k]
    br, bgg, bb = resolved["background"]
    l_bg = _luminance_srgb(br, bgg, bb)
    on_light = l_bg > _BG_LUMINANCE_LIGHT
    floor = 3.0
    for k in ("title", "body", "muted"):
        fr, fgg, fb = resolved[k]
        if _contrast_ratio_rgb((fr, fgg, fb), (br, bgg, bb)) >= floor:
            continue
        if on_light:
            dr, dg, db = _DEFAULT_DARK_ON_LIGHT[k]
        else:
            dr, dg, db = _DEFAULT_THEME_RGB[k]
        raw[k] = f"#{dr:02X}{dg:02X}{db:02X}"
    return raw


def _normalize_theme_dict(data: Any) -> dict[str, str] | None:
    """Waliduje kolory z LLM: klucze background/title/body/muted, wartości #RRGGBB."""
    if not isinstance(data, dict):
        return None
    out: dict[str, str] = {}
    for k in _THEME_JSON_KEYS:
        v = data.get(k)
        if v is None:
            continue
        s = str(v).strip()
        parsed = _parse_hex_rgb(s)
        if not parsed:
            continue
        r, g, b = parsed
        out[k] = f"#{r:02X}{g:02X}{b:02X}"
    if not out:
        return None
    _mend_theme_contrast(out)
    return out


def ensure_theme_persisted(
    new_spec: dict[str, Any], previous_spec: dict[str, Any] | None
) -> dict[str, Any]:
    """Przy edycji: gdy nowy JSON nie zawiera motywu, kopiuj z poprzedniej wersji."""
    if not previous_spec:
        return new_spec
    if _normalize_theme_dict(new_spec.get("theme")):
        return new_spec
    prev = _normalize_theme_dict(previous_spec.get("theme"))
    if not prev:
        return new_spec
    merged = {**new_spec, "theme": prev}
    return merged


def normalize_presentation_spec(data: Any) -> dict[str, Any] | None:
    """Waliduje i porządkuje słownik z modelu. Zwraca None, gdy brak slajdów merytorycznych."""
    if not isinstance(data, dict):
        return None
    title = (data.get("title") or "").strip() or "Prezentacja"
    desc = (data.get("description") or "").strip()
    slides_raw = data.get("slides")
    if not isinstance(slides_raw, list):
        return None
    out_slides: list[dict[str, Any]] = []
    for s in slides_raw:
        if not isinstance(s, dict):
            continue
        st = (s.get("title") or "").strip() or "Slajd"
        bullets = s.get("bullets")
        if not isinstance(bullets, list):
            bullets = []
        b_clean = [str(b).strip() for b in bullets if str(b).strip()]
        raw_img = s.get("image", None)
        inc = s.get("include_image")
        if inc is None:
            if raw_img is None:
                inc = False
            elif isinstance(raw_img, bool):
                inc = raw_img
            elif isinstance(raw_img, dict) and not raw_img:
                inc = False
            else:
                inc = True
        else:
            inc = bool(inc)
        image_hint: str | None = None
        if inc and isinstance(raw_img, dict):
            image_hint = (
                (raw_img.get("suggested_prompt") or raw_img.get("prompt") or "").strip() or None
            )
        elif inc and isinstance(raw_img, str) and raw_img.strip():
            image_hint = raw_img.strip()
        if not image_hint and isinstance(s.get("image_hint"), str) and s.get("image_hint", "").strip():
            image_hint = (s.get("image_hint") or "").strip()
        out_slides.append(
            {
                "title": st,
                "bullets": b_clean,
                "include_image": inc,
                "image_hint": image_hint,
            }
        )
    has_substance = bool(out_slides) or bool(desc and desc.strip()) or (title != "Prezentacja")
    if not has_substance:
        return None
    theme: dict[str, str] | None = _normalize_theme_dict(data.get("theme") if isinstance(data, dict) else None)
    out: dict[str, Any] = {
        "version": 1,
        "title": title,
        "description": desc,
        "slides": out_slides,
    }
    if theme:
        out["theme"] = theme
    return out


def parse_presentation_json(raw: str) -> dict[str, Any] | None:
    """Parsuje odpowiedź LLM (opcjonalnie w ```json)."""
    t = (raw or "").strip()
    if t.startswith("```"):
        t = t.split("\n", 1)[-1]
        if "```" in t:
            t = t.rsplit("```", 1)[0]
    try:
        data = json.loads(t)
    except (json.JSONDecodeError, TypeError):
        m = re.search(r"\{[\s\S]*\}\s*$", t)
        if m:
            try:
                data = json.loads(m.group(0))
            except json.JSONDecodeError:
                return None
        else:
            return None
    return normalize_presentation_spec(data)


def spec_to_json_text(spec: dict[str, Any]) -> str:
    return json.dumps(spec, ensure_ascii=False, indent=2)


# --- PPTX — motyw kolorowy (LLM w specie lub domyślna paleta) ---


def _resolved_theme_colors(theme: dict[str, str] | None) -> dict[str, Any]:
    from pptx.dml.color import RGBColor

    out: dict[str, Any] = {}
    thm = _normalize_theme_dict(theme) if isinstance(theme, dict) and theme else None
    for k in _THEME_JSON_KEYS:
        raw = (thm or {}).get(k) if thm else None
        tup: tuple[int, int, int] | None
        if raw:
            tup = _parse_hex_rgb(str(raw))
        else:
            tup = None
        if tup is None:
            tup = _DEFAULT_THEME_RGB[k]
        r, g, b = tup
        out[k] = RGBColor(r, g, b)
    return out


def _apply_colorful_theme_to_slide(slide: Any, colors: dict[str, Any]) -> None:
    """Tło i tekst wg palety (z LLM lub domyślnej)."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.util import Pt

    BG = colors["background"]
    TITLE = colors["title"]
    BODY = colors["body"]
    MUTED = colors["muted"]

    try:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG
    except Exception:
        return

    def _ph_type(shape: Any) -> Any:
        try:
            if shape.is_placeholder:
                return shape.placeholder_format.type
        except Exception:
            pass
        return None

    _sub_ph = (PP_PLACEHOLDER.SUBTITLE,)
    if hasattr(PP_PLACEHOLDER, "CENTER_SUBTITLE"):
        _sub_ph = (PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.CENTER_SUBTITLE)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        ph_t = _ph_type(shape)
        is_title = ph_t in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
        is_sub = ph_t in _sub_ph

        for p in tf.paragraphs:
            t = (p.text or "").strip()
            line_muted = t.startswith(("[", "("))
            if is_sub or (line_muted and not is_title):
                color = MUTED
            elif is_title:
                color = TITLE
            else:
                color = BODY
            for run in p.runs:
                run.font.color.rgb = color
            if not p.runs:
                p.font.color.rgb = color
            if is_title:
                p.font.bold = True
                try:
                    p.font.size = Pt(32)
                except Exception:
                    pass
        try:
            tf.word_wrap = True
        except Exception:
            pass


def apply_colorful_theme_to_presentation(
    prs: Any, theme: dict[str, str] | None = None
) -> None:
    """Stosuj do wszystkich slajdów po zbudowaniu treści. ``theme`` z pola ``spec.theme`` (hex #RRGGBB)."""
    color_map = _resolved_theme_colors(theme)
    for s in prs.slides:
        _apply_colorful_theme_to_slide(s, color_map)


def _cover_subtitle_font_pt(text_len: int) -> int:
    """Czcionka opisu na okładce — mniejsza przy długim tekście (w tym agendzie)."""
    if text_len < 200:
        return 18
    if text_len < 420:
        return 16
    if text_len < 800:
        return 14
    if text_len < 1200:
        return 12
    return 11


def _content_body_font_pt(num_nonempty_lines: int, max_line_len: int) -> int:
    """Czcionka listy na slajdzie treści — zmniejsz przy wielu lub długich punktach."""
    n = num_nonempty_lines
    m = max_line_len
    pt = 18
    if n > 6 or m > 100:
        pt = 14
    elif n > 4 or m > 78:
        pt = 15
    elif n > 3 or m > 58:
        pt = 16
    elif n > 2 and m > 50:
        pt = 17
    return max(11, min(pt, 20))


def _set_textframe_font_pt(tf: Any, pt: int) -> None:
    from pptx.util import Pt

    s = max(11, min(int(pt), 24))
    for p in tf.paragraphs:
        p.font.size = Pt(s)
        for r in p.runs:
            r.font.size = Pt(s)


def spec_to_pptx_bytes(
    spec: dict[str, Any],
    *,
    slide_images: dict[int, bytes] | None = None,
) -> bytes:
    """
    `slide_images`: indeks (0 = pierwszy slajd merytoryczny) → bajty PNG/JPEG
    do osadzenia po prawej; brak wpisu = sam tekst (ew. wiersz „[Propozycja grafiki: …]”).
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title = (spec.get("title") or "Prezentacja")[:200]
    desc = (spec.get("description") or "")[:2200]
    # Slajd 1 — tytuł + opis
    layout0 = prs.slide_layouts[0]
    slide0 = prs.slides.add_slide(layout0)
    slide0.shapes.title.text = title
    if len(slide0.placeholders) > 1:
        ph = slide0.placeholders[1]
        dplain = desc if desc else " "
        ph.text = dplain
        try:
            _set_textframe_font_pt(
                ph.text_frame, _cover_subtitle_font_pt(len(dplain))
            )
        except Exception:
            pass

    layout1 = prs.slide_layouts[1]
    sim = slide_images or {}
    for slide_idx, s in enumerate(spec.get("slides") or []):
        if not isinstance(s, dict):
            continue
        st = (s.get("title") or "Slajd")[:200]
        bullets = s.get("bullets") or []
        body_lines: list[str] = list(bullets) if isinstance(bullets, list) else []
        embed_bytes: bytes | None = sim.get(slide_idx) or None
        if s.get("include_image") and not embed_bytes:
            hint = (s.get("image_hint") or "").strip()
            if hint:
                body_lines.append("")
                body_lines.append(f"[Propozycja grafiki: {hint}]")
            else:
                body_lines.append("")
                body_lines.append("(Miejsce na grafikę — uzupełnij w PowerPoint, jeśli potrzeba.)")
        slide = prs.slides.add_slide(layout1)
        slide.shapes.title.text = st
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        nonempty = [x for x in body_lines if str(x).strip()]
        mlen = max((len(x) for x in nonempty), default=0)
        base_pt = _content_body_font_pt(len(nonempty), mlen)
        if not body_lines:
            p0 = tf.paragraphs[0]
            p0.text = " "
            p0.font.size = Pt(base_pt)
        for i, line in enumerate(body_lines):
            small = line.startswith("[") or line.startswith("(")
            fs = max(11, base_pt - 1) if small else base_pt
            if i == 0:
                p = tf.paragraphs[0]
                p.text = line
                p.font.size = Pt(fs)
            else:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(fs)
                p.level = 0
        if embed_bytes:
            try:
                body.left = Inches(0.4)
                body.top = Inches(1.2)
                body.width = Inches(5.7)
                body.height = Inches(4.7)
                stream = io.BytesIO(embed_bytes)
                slide.shapes.add_picture(stream, Inches(6.1), Inches(1.2), height=Inches(4.65))
            except Exception:
                pass

    th = spec.get("theme")
    spec_theme = th if isinstance(th, dict) else None
    apply_colorful_theme_to_presentation(prs, spec_theme)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _shape_text(sh) -> str:
    try:
        if hasattr(sh, "text") and sh.text:
            return sh.text.strip()
    except Exception:
        pass
    return ""


def pptx_to_spec(data: bytes) -> dict[str, Any] | None:
    """Odczyt istniejącego PPTX do specyfikacji (edycja / kontynuacja)."""
    try:
        from pptx import Presentation
    except Exception:
        return None
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception:
        return None
    if not prs.slides:
        return None
    s0 = prs.slides[0]
    title = _shape_text(s0.shapes.title) if s0.shapes.title else "Prezentacja"
    desc = ""
    if len(s0.shapes) > 1:
        for sh in s0.shapes:
            if sh == s0.shapes.title:
                continue
            t = _shape_text(sh)
            if t and t != title:
                desc = t
                break
    slides_out: list[dict[str, Any]] = []
    for idx in range(1, len(prs.slides)):
        sl = prs.slides[idx]
        st = _shape_text(sl.shapes.title) if sl.shapes.title else f"Slajd {idx + 1}"
        body_t = ""
        if len(sl.placeholders) > 1:
            try:
                body_t = _shape_text(sl.placeholders[1])
            except Exception:
                body_t = ""
        lines: list[str] = []
        include_image = bool(body_t) and (
            "[Propozycja grafiki:" in body_t or "Miejsce na grafikę" in body_t
        )
        image_hint: str | None = None
        if "[Propozycja grafiki:" in body_t and "]" in body_t:
            a = body_t.find("[Propozycja grafiki:")
            b = body_t.rfind("]")
            if a >= 0 and b > a:
                image_hint = body_t[a + len("[Propozycja grafiki:") : b].strip() or None
        for part in body_t.split("\n"):
            p = part.strip()
            if not p:
                continue
            if p.startswith("[Propozycja grafiki:") and p.endswith("]"):
                continue
            if p.startswith("(Miejsce na grafikę"):
                continue
            lines.append(p)
        slides_out.append(
            {
                "title": st,
                "bullets": lines,
                "include_image": include_image,
                "image_hint": image_hint,
            }
        )
    return normalize_presentation_spec(
        {"title": title, "description": desc, "slides": slides_out},
    )


def spec_to_readable_plan_text(spec: dict[str, Any]) -> str:
    """Czytelny opis planu (tytuł, opis, slajdy) — do pliku PDF i podglądu."""
    t_raw = _normalize_theme_dict(spec.get("theme")) if spec.get("theme") else None
    theme_line = ""
    if t_raw:
        theme_line = "Motyw kolorów: " + ", ".join(f"{k}={t_raw[k]}" for k in _THEME_JSON_KEYS if k in t_raw) + "\n\n"
    parts: list[str] = [f"# {spec.get('title', '')}", "", (spec.get("description") or "").strip(), "", theme_line]
    for i, s in enumerate(spec.get("slides") or [], start=1):
        st = s.get("title", "") if isinstance(s, dict) else ""
        parts.append(f"## Slajd {i + 1} — {st}")
        bullets = (s.get("bullets") or []) if isinstance(s, dict) else []
        if isinstance(bullets, list):
            for b in bullets:
                parts.append(f"- {b}")
        if isinstance(s, dict) and s.get("include_image"):
            h = (s.get("image_hint") or "").strip()
            if h:
                parts.append(f"  (propozycja grafiki: {h})")
            else:
                parts.append("  (miejsce na grafikę — opcjonalnie)")
        parts.append("")
    return "\n".join(parts).strip()


def extract_pptx_plain_text(data: bytes) -> str:
    """Treść PPTX do indeksu wyszukiwania / eksportu do PDF (tekstowo)."""
    spec = pptx_to_spec(data)
    if not spec:
        return ""
    return spec_to_readable_plan_text(spec)
